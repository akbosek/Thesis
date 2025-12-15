from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()

# Helper function to create slides
def create_slide(prs, title_text, content_items, placeholder_text="[PASTE GRAPHS HERE]"):
    # Use 'Title and Content' layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Set Title
    title = slide.shapes.title
    title.text = title_text
    
    # Adjust text box (Left side)
    body = slide.placeholders[1]
    body.left = Inches(0.5)
    body.top = Inches(1.5)
    body.width = Inches(5.5)
    body.height = Inches(5.5)
    
    tf = body.text_frame
    tf.clear()
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        p.space_after = Pt(10)
        if 'bold' in item:
            p.font.bold = True
        if 'size' in item:
            p.font.size = Pt(item['size'])
        else:
            p.font.size = Pt(16)

    # Placeholder for Graphs (Right side)
    left = Inches(6.2)
    top = Inches(1.5)
    width = Inches(3.5)
    height = Inches(4.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf_box = textbox.text_frame
    tf_box.text = placeholder_text
    tf_box.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_box.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    
    # Add a border rectangle to indicate paste area
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.background() # Transparent fill
    line = shape.line
    line.color.rgb = RGBColor(200, 200, 200)
    line.width = Pt(1.5)

# --- SLIDE 1: TITLE ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Bitcoin Price Prediction Using Hybrid CNN-LSTM Ensemble"
subtitle.text = "Optimizing Trading Strategies through Macro-Financial Data and Threshold Calibration\n\nProject Presentation"

# --- SLIDE 2: THE PROBLEM & SOLUTION ---
content = [
    {"text": "The Challenge: 'Long Bias' & Noise", "level": 0, "bold": True},
    {"text": "Problem: Bitcoin has a historical uptrend. Standard AI models become 'optimistic' and fail to predict crashes.", "level": 1},
    {"text": "Problem: Crypto markets contain high noise. Models struggle to distinguish signals from random fluctuations.", "level": 1},
    {"text": "Our Solution:", "level": 0, "bold": True},
    {"text": "1. Hybrid Architecture (CNN-LSTM): Combining pattern recognition with trend following.", "level": 1},
    {"text": "2. Macro Data Injection: Using S&P500, Oil, and VIX to provide context.", "level": 1},
    {"text": "3. Advanced Calibration: Asymmetric thresholds and class weights to force 'Short' predictions.", "level": 1}
]
create_slide(prs, "Project Overview", content, "[PASTE ARCHITECTURE DIAGRAM OR DATA HEAD]")

# --- SLIDE 3: METHODOLOGY - ARCHITECTURE ---
content = [
    {"text": "Engine: CNN-LSTM Hybrid", "level": 0, "bold": True},
    {"text": "Based on: Somayajulu et al. (2024)", "level": 1},
    {"text": "WHAT: We combined two neural networks.", "level": 1},
    {"text": "- CNN (Convolutional): Acts as a 'Technical Analyst'. It scans candle shapes and local patterns, filtering out noise.", "level": 2},
    {"text": "- LSTM (Long Short-Term Memory): Acts as a 'Strategic Investor'. It remembers 30-day sequences to identify broader trends.", "level": 2},
    {"text": "EFFECT: The model is resistant to sudden 'wicks' and price spikes, focusing on sustained momentum.", "level": 1},
    {"text": "Ensemble Learning:", "level": 0, "bold": True},
    {"text": "We train 5 independent models and average their votes. This 'Wisdom of the Crowd' reduces variance and risk.", "level": 1}
]
create_slide(prs, "Methodology: The 'Engine'", content, "[PASTE MODEL SUMMARY OR CODE SNIPPET]")

# --- SLIDE 4: METHODOLOGY - DATA ---
content = [
    {"text": "Fuel: Macro-Financial Data", "level": 0, "bold": True},
    {"text": "Based on: Parim et al. (2024)", "level": 1},
    {"text": "WHY: Bitcoin does not exist in a vacuum. It correlates with global liquidity.", "level": 1},
    {"text": "Added Features:", "level": 1},
    {"text": "- VIX (Fear Index): Measures market panic. High VIX often precedes crypto dumps.", "level": 2},
    {"text": "- S&P 500 & NASDAQ: Captures correlation with traditional equities.", "level": 2},
    {"text": "- Crude Oil (WTI): Proxy for inflation and global energy costs.", "level": 2},
    {"text": "EFFECT: The model can anticipate moves based on external shocks, not just past prices.", "level": 1}
]
create_slide(prs, "Methodology: The 'Fuel'", content, "[PASTE DATA CORRELATION MATRIX OR HEAD]")

# --- SLIDE 5: EXPERIMENT A (OPTIMAL) ---
content = [
    {"text": "Experiment A: The 'Golden Mean' (Optimal)", "level": 0, "bold": True},
    {"text": "Configuration:", "level": 1},
    {"text": "- Manual Thresholds: Long > 0.55 | Short < 0.48", "level": 2},
    {"text": "- Class Weights: 1.3 (Balanced penalty)", "level": 2},
    {"text": "Results:", "level": 1},
    {"text": "- Validation Win Rate: ~58.75%", "level": 2},
    {"text": "- Sharpe Ratio: 0.79", "level": 2},
    {"text": "- Action Rate: ~11% (80 trades)", "level": 2},
    {"text": "Analysis:", "level": 1},
    {"text": "This configuration offers the best balance. The Equity Curve is smooth, and the model trades frequently enough to be profitable without over-trading.", "level": 2}
]
create_slide(prs, "Exp A: The Optimal Configuration", content, "[PASTE GRAPHS FROM 'DOBRY' FOLDER:\nEquity, Signals, Hist, Matrix]")

# --- SLIDE 6: EXPERIMENT B (SNIPER) ---
content = [
    {"text": "Experiment B: The 'Sniper' Approach", "level": 0, "bold": True},
    {"text": "Configuration:", "level": 1},
    {"text": "- Wide Dead Zone (Very strict confidence required)", "level": 2},
    {"text": "- Strong L2 Regularization (0.0001)", "level": 2},
    {"text": "Results:", "level": 1},
    {"text": "- Validation Win Rate: 60.00% (Highest Accuracy)", "level": 2},
    {"text": "- Action Rate: 5.71% (Only ~40 trades)", "level": 2},
    {"text": "Analysis:", "level": 1},
    {"text": "High precision, but too passive. Proves the model can identify high-probability setups, but misses many opportunities.", "level": 2}
]
create_slide(prs, "Exp B: The 'Sniper' Approach", content, "[PASTE GRAPHS FROM '60%' FOLDER]")

# --- SLIDE 7: EXPERIMENT C (SHORTING) ---
content = [
    {"text": "Experiment C: Stress Test (Forcing Shorts)", "level": 0, "bold": True},
    {"text": "Configuration:", "level": 1},
    {"text": "- Aggressive Class Weights (Boosted Short importance)", "level": 2},
    {"text": "- Lowered Thresholds to force activity", "level": 2},
    {"text": "Results:", "level": 1},
    {"text": "- Action Rate: 58.29% (Over-trading)", "level": 2},
    {"text": "- Win Rate: 52.70% (Drop in quality)", "level": 2},
    {"text": "Analysis:", "level": 1},
    {"text": "Proves we CAN force the model to short against the trend, but 'Quantity != Quality'. Fighting the trend reduces the Sharpe Ratio.", "level": 2}
]
create_slide(prs, "Exp C: Forcing Shorts", content, "[PASTE GRAPHS FROM 'SHORTY' FOLDER]")

# --- SLIDE 8: EXPERIMENT D (STABILITY) ---
content = [
    {"text": "Experiment D: Auto-Calibrated Baseline", "level": 0, "bold": True},
    {"text": "Configuration:", "level": 1},
    {"text": "- Dynamic Median Calibration (Model finds its own zero-point)", "level": 2},
    {"text": "Results:", "level": 1},
    {"text": "- Win Rate: 56.32%", "level": 2},
    {"text": "- Sharpe Ratio: 0.94 (Excellent Risk/Reward)", "level": 2},
    {"text": "Analysis:", "level": 1},
    {"text": "Zero overfitting (Train ~ Val). The safest 'set-and-forget' strategy.", "level": 2}
]
create_slide(prs, "Exp D: Auto-Calibrated Stability", content, "[PASTE GRAPHS FROM '56%' FOLDER]")

# --- SLIDE 9: SOLVING THE 'LONG BIAS' ---
content = [
    {"text": "The Problem:", "level": 0, "bold": True},
    {"text": "Neural Networks are optimistic about Bitcoin. Median prediction > 0.50. A standard threshold results in ZERO Shorts.", "level": 1},
    {"text": "The Solution: Asymmetric Thresholding", "level": 0, "bold": True},
    {"text": "We decoupled the decision logic:", "level": 1},
    {"text": "- LONG Condition: Confidence > 0.55 (Harder)", "level": 2},
    {"text": "- SHORT Condition: Confidence < 0.48 (Easier)", "level": 2},
    {"text": "Effect:", "level": 1},
    {"text": "We 'tricked' the model's optimism. By lowering the Short bar, we captured downturns even when the model wasn't 100% bearish.", "level": 2}
]
create_slide(prs, "Key Breakthrough: Solving Long Bias", content, "[PASTE HISTOGRAM SHOWING SHIFT]")

# --- SLIDE 10: CONCLUSION ---
content = [
    {"text": "Summary:", "level": 0, "bold": True},
    {"text": "Successfully built a profitable AI trading bot (Sharpe > 0.80) that outperforms random chance.", "level": 1},
    {"text": "Key Findings:", "level": 0, "bold": True},
    {"text": "1. Macro Data Matters: VIX and SP500 improved signal quality [Parim et al.].", "level": 1},
    {"text": "2. Logic > Architecture: Tuning thresholds had a bigger impact than adding neurons.", "level": 1},
    {"text": "3. The 'Golden Mean' configuration provides the most viable live-trading strategy.", "level": 1}
]
create_slide(prs, "Conclusion", content, "[PASTE SUMMARY TABLE OR GRAPH]")

# Save
file_name = "Bitcoin_Trading_Bot_Presentation.pptx"
prs.save(file_name)
print(f"✅ Presentation saved as: {file_name}")