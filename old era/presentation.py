from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Tworzenie nowej prezentacji
prs = Presentation()

# Funkcja pomocnicza do dodawania slajdów
def add_slide(prs, title_text, content_text_list, placeholder_text="[PASTE YOUR GRAPHS HERE]"):
    # Wybór układu slajdu (Title and Content)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Ustawienie tytułu
    title = slide.shapes.title
    title.text = title_text
    
    # Ustawienie treści (lewa strona)
    body = slide.placeholders[1]
    body.width = Inches(5.5) 
    body.height = Inches(6.0)
    
    tf = body.text_frame
    tf.clear() 
    
    for item in content_text_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16) # Trochę mniejsza czcionka, żeby zmieścić więcej treści
        p.space_after = Pt(8)

    # Ramka na wykresy
    left = Inches(6.0)
    top = Inches(1.5)
    width = Inches(3.5)
    height = Inches(4.0)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf_box = textbox.text_frame
    tf_box.text = placeholder_text
    tf_box.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# --- SLAJD 1: TYTUŁOWY ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Bitcoin Price Prediction Using Hybrid CNN-LSTM Ensemble"
subtitle.text = "Optimizing Trading Strategies through Macro-Financial Data and Threshold Calibration\n\nAuthor: [Twoje Imię]"

# --- SLAJD 2: METHODOLOGY ---
content = [
    "Based on Research Papers:",
    "1. Parim et al. (2024): Demonstrated that Bitcoin is driven by macro factors.",
    "   - Implemented: VIX (Fear Index), S&P 500, NASDAQ, Oil added to dataset.",
    "2. Somayajulu et al. (2024): Proposed CNN-LSTM Hybrid.",
    "   - Implemented: 1D Convolution (feature extraction) -> LSTM (sequence).",
    "   - Adopted 'Adam' optimizer and Dropout layers for stability."
]
add_slide(prs, "Methodology & Theoretical Foundation", content, "[PASTE ARCHITECTURE DIAGRAM OR DATA HEAD]")

# --- SLAJD 3: EXPERIMENT A (OPTIMAL) ---
content = [
    "Experiment A: The 'Golden Mean' (Optimal)",
    "Based on file: dobry.txt",
    "",
    "Configuration:",
    "- Manual Thresholds: Long > 0.55 | Short < 0.48",
    "- Class Weights: 1.3 (Balanced approach)",
    "",
    "Results:",
    "- Validation Win Rate: 58.75%",
    "- Sharpe Ratio: 0.79",
    "- Action Rate: 11.43% (80 trades)",
    "",
    "Conclusion: Best balance between profit stability and risk."
]
add_slide(prs, "Exp A: The 'Golden Mean' (Optimal)", content, "[PASTE 4 GRAPHS FROM 'DOBRY' FOLDER]")

# --- SLAJD 4: EXPERIMENT B (SNIPER) ---
content = [
    "Experiment B: The 'Sniper' Approach",
    "Based on file: 60% wr.txt",
    "",
    "Configuration:",
    "- Wide Dead Zone (Strict confidence required)",
    "- Strong L2 Regularization (0.0001)",
    "",
    "Results:",
    "- Validation Win Rate: 60.00% (Highest Accuracy)",
    "- Action Rate: 5.71% (Only ~40 trades)",
    "",
    "Analysis:",
    "- Extremely precise but too passive for active trading.",
    "- Proves the model can identify high-probability setups."
]
add_slide(prs, "Exp B: The 'Sniper' Approach (Precision)", content, "[PASTE GRAPHS FROM '60%' FOLDER]")

# --- SLAJD 5: EXPERIMENT C (FORCING SHORTS) ---
content = [
    "Experiment C: Forcing Short Positions",
    "Based on file: shorty.txt",
    "",
    "Configuration:",
    "- Aggressive Class Weights (Boosted Short importance)",
    "- Lowered Thresholds to force activity",
    "",
    "Results:",
    "- Action Rate: 58.29% (High Activity / Over-trading)",
    "- Validation Win Rate: 52.70%",
    "- Sharpe Ratio: 0.38 (High volatility)",
    "",
    "Analysis:",
    "- Proves we CAN force the model to short against the trend.",
    "- However, Quantity != Quality. Fighting the trend reduces Sharpe."
]
add_slide(prs, "Exp C: Forcing Shorts (Stress Test)", content, "[PASTE GRAPHS FROM 'SHORTY' FOLDER]")

# --- SLAJD 6: EXPERIMENT D (STABILITY) ---
content = [
    "Experiment D: Auto-Calibrated Baseline",
    "Based on file: 56%.txt",
    "",
    "Configuration:",
    "- Dynamic Median Calibration (Model finds its own zero-point)",
    "- Dropout: 0.3",
    "",
    "Results:",
    "- Validation Win Rate: 56.32%",
    "- Sharpe Ratio: 0.94 (Excellent Risk/Reward)",
    "",
    "Conclusion: Safe, set-and-forget configuration with no overfitting."
]
add_slide(prs, "Exp D: Auto-Calibrated (Stability)", content, "[PASTE GRAPHS FROM '56%' FOLDER]")

# --- SLAJD 7: THE SHORT CHALLENGE ---
content = [
    "The Problem: 'Long Bias'",
    "- Bitcoin has a historical upward trend (2014-2024).",
    "- Neural Networks naturally become optimistic (Median Prediction > 0.50).",
    "- Standard threshold (0.50) results in ZERO Short positions.",
    "",
    "Our Solution:",
    "- Asymmetric Thresholding (Shifted Logic).",
    "- We manually lowered the Short Threshold (e.g., to <0.48).",
    "- This forced the model to act on 'relative weakness' rather than absolute downtrends."
]
add_slide(prs, "Challenge: Overcoming Long Bias", content, "[PASTE HISTOGRAM SHOWING SHIFT]")

# --- SLAJD 8: CONCLUSION ---
content = [
    "Summary:",
    "- Built a Hybrid CNN-LSTM Ensemble that beats random chance (~58% vs 50%).",
    "- Validated Parim et al.'s thesis on Macro Data importance.",
    "",
    "Key Findings:",
    "- Hyperparameter tuning of Thresholds was more impactful than architecture changes.",
    "- 'Class Weights' and 'Asymmetric Thresholds' solved the Shorting problem.",
    "",
    "Future Work:",
    "- Implement 'Dual-Core' architecture (Separate Bull/Bear models).",
    "- Add On-Chain data (Hash Rate, Wallets)."
]
add_slide(prs, "Conclusion & Future Path", content, "[PASTE FINAL SUMMARY GRAPH]")

# Zapisanie pliku
output_file = "Bitcoin_Project_Presentation_Final.pptx"
prs.save(output_file)
print(f"✅ Prezentacja gotowa: {output_file}")